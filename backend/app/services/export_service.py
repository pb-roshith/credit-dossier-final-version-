"""
Export Service — generates PPT, PDF, and DOCX documents from deal narratives.
Also provides a combined PDF report from all sections.
"""

from __future__ import annotations

import io
import logging
import re
import base64
import html
from datetime import datetime

import markdown
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from sqlalchemy.orm import Session, joinedload
from reportlab import rl_config

# CVE-2020-28463: PDF exports only need locally generated data: images. A
# non-routable sentinel keeps ReportLab's host allowlist enabled while the
# scheme allowlist blocks HTTP(S), file, FTP, and other external resources.
rl_config.trustedHosts = ["no-remote-resources.invalid"]
rl_config.trustedSchemes = ["data"]

from xhtml2pdf import pisa
from html2docx import html2docx

from app.models.deal import Deal, Section
from app.security import mask_sensitive_text

logger = logging.getLogger(__name__)


def _format_amount(amount: float, currency: str) -> str:
    if currency == "INR":
        cr = amount / 10_000_000
        return f"INR {cr:,.2f} Cr"
    return f"{currency} {amount:,.0f}"


def _parse_table_data(table_str: str) -> pd.DataFrame | None:
    lines = table_str.strip().split('\n')
    if len(lines) < 3: 
        return None
    
    headers = [col.replace('**', '').replace('*', '').strip() for col in lines[0].split('|')[1:-1]]
    rows = []
    for line in lines[2:]:
        cols = [col.replace('**', '').replace('*', '').strip() for col in line.split('|')[1:-1]]
        if len(cols) == len(headers):
            rows.append(cols)
    if not rows: 
        return None
    
    df = pd.DataFrame(rows, columns=headers)
    
    for col in df.columns[1:]:
        df[col] = pd.to_numeric(df[col].str.replace(r'[^0-9.-]', '', regex=True), errors='coerce')
    
    if df.shape[1] > 1 and df.iloc[:, 1:].notna().any().any():
        return df
    return None


def _clean_unprintable_chars(text: str) -> str:
    """Removes emojis, dingbats, and unsupported unicode characters that break PDF rendering."""
    # Keep ASCII, extended Latin, common punctuation, Euro, Pound, and dashes
    # This regex removes characters outside of these safe ranges.
    clean_text = re.sub(r'[^\x00-\x7F\xA0-\xFF\u0100-\u017F\u2013-\u2014\u2018-\u201D\u20AC\u00A3]', '', text)
    return clean_text


def _section_narrative(section: Section) -> str:
    """Return final narrative text with UI-only source markers removed."""
    content = section.final_generated_content or section.generated_content or ""
    # Current inline citation format.
    content = re.sub(
        r"\s*\[Source\s*:\s*[^\]\r\n]+\]",
        "",
        content,
        flags=re.IGNORECASE,
    )
    # Remove legacy numbered reference sections and their inline markers.
    content = re.sub(
        r"\n{0,2}#{1,6}\s*(?:References|Sources|Bibliography)\b.*\Z",
        "",
        content,
        flags=re.IGNORECASE | re.DOTALL,
    )
    content = re.sub(r"\[(?:\d{1,2})(?:\s*,\s*\d{1,2})*\]", "", content)
    return mask_sensitive_text(content.strip())

def _markdown_to_html_with_graphs(text: str, theme_palette: list[str]) -> str:
    """Detects numeric markdown tables, injects matplotlib charts, and converts to HTML."""
    # Preserve Markdown while rendering any embedded HTML inert. Locally
    # generated chart tags are inserted only after this escaping step.
    text = html.escape(text, quote=False)
    table_pattern = re.compile(r'(^[ \t]*\|[^\n]+\|[ \t]*\n[ \t]*\|[-:| ]+\|[ \t]*\n(?:[ \t]*\|[^\n]+\|[ \t]*(?:\n|$))+)', re.MULTILINE)
    
    def replacer(match):
        table_str = match.group(1)
        # Ensure table has empty lines around it so python-markdown parses it correctly
        formatted_table = f"\n\n{table_str.strip()}\n\n"
        
        df = _parse_table_data(table_str)
        if df is not None:
            try:
                fig, ax = plt.subplots(figsize=(10, 6))
                # Use the dynamic colors from the full palette
                colors = theme_palette if len(theme_palette) >= 5 else (theme_palette + ["#64748b", "#94a3b8", "#cbd5e1", "#334155", "#0f172a"])
                bars = df.plot(x=df.columns[0], kind='bar', ax=ax, color=colors[:len(df.columns)-1], width=0.7)
                
                def format_label(val):
                    if pd.isna(val) or val == 0:
                        return ""
                    abs_val = abs(val)
                    if abs_val >= 1_000_000:
                        return f"{val/1_000_000:.1f}M"
                    elif abs_val >= 1_000:
                        return f"{val/1_000:.1f}K"
                    elif abs_val >= 1:
                        return f"{val:.1f}"
                    else:
                        return f"{val:.2f}"

                # Add value labels
                for container in ax.containers:
                    labels = [format_label(v.get_height()) for v in container]
                    ax.bar_label(container, labels=labels, padding=3, color='#334155', fontsize=9, fontweight='bold', rotation=90)
                
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
                ax.spines['left'].set_visible(False)
                ax.set_xlabel("")
                ax.set_yticks([]) # Hide y-ticks to rely on bar labels
                plt.xticks(rotation=45, ha='right', fontsize=10)
                ax.legend(loc='upper center', bbox_to_anchor=(0.5, 1.15), ncol=min(3, len(df.columns)-1), frameon=False)
                plt.tight_layout()
                
                buf = io.BytesIO()
                fig.savefig(buf, format='png', dpi=300, transparent=True, bbox_inches='tight')
                plt.close(fig)
                
                img_b64 = base64.b64encode(buf.getvalue()).decode('utf-8')
                img_tag = f"<div style='text-align:center; margin: 20px 0;'><img src='data:image/png;base64,{img_b64}' style='max-width: 100%; height: auto;' /></div>\n\n"
                return formatted_table + img_tag
            except Exception as e:
                logger.error(f"Failed to plot table: {e}")
                return formatted_table
        return formatted_table

    transformed = table_pattern.sub(replacer, text)
    rendered_html = markdown.markdown(transformed, extensions=['tables'])
    
    # Post-process HTML to force equal column widths in PDF
    def fix_table_widths(match):
        table_html = match.group(0)
        # Find number of columns from the first row
        tr_match = re.search(r'<tr\b[^>]*>(.*?)</tr>', table_html, re.DOTALL | re.IGNORECASE)
        if not tr_match:
            return table_html
            
        cells = re.findall(r'<(?:th|td)\b', tr_match.group(1), re.IGNORECASE)
        if not cells:
            return table_html
            
        col_count = len(cells)
        if col_count == 0:
            return table_html
            
        width_pct = int(100 / col_count)
        
        # Inject HTML width attribute into ALL th and td tags to force xhtml2pdf to obey boundaries
        table_html = re.sub(r'<(th|td)\b([^>]*)>', f'<\\1 width="{width_pct}%" \\2>', table_html, flags=re.IGNORECASE)
        return table_html

    rendered_html = re.sub(
        r'<table\b.*?>.*?</table>',
        fix_table_widths,
        rendered_html,
        flags=re.DOTALL | re.IGNORECASE,
    )
    
    return rendered_html


class ExportService:
    """Generates export documents from deal data."""

    @staticmethod
    def get_deal_with_sections(db: Session, deal_id: str) -> Deal | None:
        return (
            db.query(Deal)
            .options(joinedload(Deal.sections), joinedload(Deal.versions))
            .filter(Deal.id == deal_id)
            .first()
        )

    @staticmethod
    def _generate_html_document(deal: Deal, for_pdf: bool = True) -> str:
        # Filter generated sections
        valid_sections = [
            s for s in sorted(deal.sections, key=lambda x: x.order_index) 
            if s.state == "ready" and _section_narrative(s).strip()
        ]

        def safe_color(value: object, fallback: str) -> str:
            candidate = str(value or "")
            return candidate if re.fullmatch(r"#[0-9A-Fa-f]{6}", candidate) else fallback

        def safe_text(value: object) -> str:
            return html.escape(mask_sensitive_text(str(value or "")))

        p_color = safe_color(getattr(deal, "primary_color", None), "#002060")
        s_color = safe_color(getattr(deal, "secondary_color", None), "#800020")
        
        # Parse full palette
        import json
        raw_palette = getattr(deal, "theme_palette", None)
        if raw_palette:
            try:
                parsed_palette = json.loads(raw_palette)
                theme_palette = [safe_color(color, "#64748b") for color in parsed_palette[:8]]
            except Exception:
                theme_palette = [p_color, s_color, "#1e293b", "#3b82f6", "#f59e0b"]
        else:
            theme_palette = [p_color, s_color, "#1e293b", "#3b82f6", "#f59e0b"]

        html_parts = []
        
        # Cover
        html_parts.append(f"<h1 style='text-align: center; font-size: 32pt; font-weight: bold; margin-top: 200px; color: {p_color};'>CREDIT PITCH BOOK</h1>")
        html_parts.append(f"<h2 style='text-align: center; font-size: 22pt; margin-top: 10px; color: {s_color};'>{safe_text(deal.customer)}</h2>")
        html_parts.append(f"<p style='text-align: center; color: #64748b; font-size: 14pt; margin-top: 30px;'>{safe_text(deal.industry)} | {safe_text(deal.segment)} | {safe_text(deal.geography)}</p>")
        html_parts.append(f"<p style='text-align: center; color: #64748b; font-size: 14pt;'>Facility: {safe_text(deal.facility)} — {safe_text(_format_amount(deal.amount, deal.currency))}</p>")
        html_parts.append(f"<p style='text-align: center; color: #94a3b8; font-size: 11pt; margin-top: 80px;'>Generated: {datetime.now().strftime('%B %d, %Y')}</p>")
        html_parts.append("<p style='text-align: center; color: #991b1b; font-size: 10pt;'><strong>CONFIDENTIAL — INTERNAL USE ONLY</strong></p>")
        
        page_break = "<pdf:nextpage />" if for_pdf else "<br/><br/><br/>"
        html_parts.append(page_break)
        
        # TOC
        html_parts.append(f"<h2 style='color: {p_color}; border-bottom: 2px solid #e2e8f0; padding-bottom: 10px;'>TABLE OF CONTENTS</h2><ul style='list-style-type: none; padding-left: 0;'>")
        for i, section in enumerate(valid_sections, 1):
            html_parts.append(f"<li style='margin-bottom: 12px; font-size: 12pt; color: #334155;'><strong style='color: {p_color};'>{i}.</strong> {safe_text(section.title)}</li>")
        html_parts.append("</ul>")
        html_parts.append(page_break)

        # Summary Table
        html_parts.append(f"<h2 style='color: {p_color}; border-bottom: 2px solid #e2e8f0; padding-bottom: 10px;'>DEAL SUMMARY</h2>")
        html_parts.append("<table border='1' cellpadding='10' style='width: 100%; border-collapse: collapse; margin-top: 20px;'>")
        html_parts.append(f"<tr><th style='background-color: {p_color}; color: white; width: 35%; font-weight: bold;'>Parameter</th><th style='background-color: {p_color}; color: white; font-weight: bold;'>Details</th></tr>")
        
        summary_items = [
            ("Customer", deal.customer),
            ("Customer Type", deal.customer_type),
            ("Industry / Sector", f"{deal.industry} / {deal.sector}"),
            ("Segment", deal.segment),
            ("Geography", deal.geography),
            ("KYC Status", deal.kyc.title()),
            ("Facility Type", deal.facility),
            ("Amount", _format_amount(deal.amount, deal.currency)),
            ("Tenure", f"{deal.tenure} months"),
            ("Pricing", deal.pricing),
            ("Repayment", deal.repayment),
            ("Collateral", "Secured" if deal.collateral else "Clean / Unsecured"),
            ("Target Date", deal.due),
        ]
        
        for k, v in summary_items:
            html_parts.append(f"<tr><td style='background-color: #f8fafc;'><strong style='color: #334155;'>{safe_text(k)}</strong></td><td>{safe_text(v)}</td></tr>")
        html_parts.append("</table>")
        html_parts.append(page_break)

        # Sections
        for i, section in enumerate(valid_sections, 1):
            html_parts.append(f"<h2 style='color: {p_color}; border-bottom: 2px solid #e2e8f0; padding-bottom: 8px;'>{i}. {safe_text(section.title)}</h2>")
            clean_content = _clean_unprintable_chars(_section_narrative(section))
            content_html = _markdown_to_html_with_graphs(clean_content, theme_palette)
            html_parts.append(content_html)
            if i < len(valid_sections):
                html_parts.append(page_break)

        css = f"""
        @page {{ size: a4 portrait; margin: 2.0cm; }}
        body {{ font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif; font-size: 11pt; line-height: 1.5; color: #334155; }}
        h1, h2, h3, h4 {{ color: {p_color}; margin-top: 25px; margin-bottom: 12px; font-weight: 600; }}
        p {{ margin-bottom: 15px; text-align: justify; }}
        table {{ width: 100%; table-layout: fixed; border-collapse: collapse; margin-top: 15px; margin-bottom: 25px; border: 1px solid #cbd5e1; -pdf-keep-with-next: false; }}
        th {{ background-color: {p_color}; color: white; padding: 8px 10px; text-align: left; font-weight: bold; border: 1px solid {p_color}; word-wrap: break-word; }}
        td {{ border: 1px solid #cbd5e1; padding: 8px 10px; vertical-align: top; word-wrap: break-word; }}
        tr:nth-child(even) {{ background-color: #f8fafc; }}
        ul, ol {{ margin-top: 10px; margin-bottom: 15px; padding-left: 20px; }}
        li {{ margin-bottom: 6px; }}
        img {{ max-width: 100%; }}
        blockquote {{ border-left: 4px solid #cbd5e1; margin-left: 0; padding-left: 15px; color: #64748b; font-style: italic; }}
        """

        full_html = f"""
        <html>
        <head>
        <style>{css}</style>
        </head>
        <body>
        {''.join(html_parts)}
        </body>
        </html>
        """
        return full_html


    @staticmethod
    def generate_pdf(deal: Deal) -> bytes:
        """Generate a perfectly formatted PDF with embedded graphs using xhtml2pdf."""
        html = ExportService._generate_html_document(deal, for_pdf=True)
        buffer = io.BytesIO()
        
        # xhtml2pdf logging is noisy, suppress it temporarily if desired, but default is fine
        pisa_status = pisa.CreatePDF(html, dest=buffer)
        
        if pisa_status.err:
            logger.error("Error generating PDF via xhtml2pdf")
            
        return buffer.getvalue()

    @staticmethod
    def generate_docx(deal: Deal) -> bytes:
        """Generate a formatted DOCX with embedded graphs using html2docx."""
        html = ExportService._generate_html_document(deal, for_pdf=False)
        buffer = html2docx(html, title="Confidential Credit Pitch Book")
        from docx import Document
        document = Document(buffer)
        document.core_properties.author = ""
        document.core_properties.last_modified_by = ""
        document.core_properties.comments = "CONFIDENTIAL - INTERNAL USE ONLY"
        secured = io.BytesIO()
        document.save(secured)
        return secured.getvalue()

    @staticmethod
    async def generate_pptx(deal: Deal) -> bytes:
        """Generate a local-only PPTX without disclosing report data externally."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import asyncio
        import re

        prs = Presentation()
        prs.core_properties.title = "Confidential Credit Pitch Book"
        prs.core_properties.author = ""
        prs.core_properties.last_modified_by = ""
        prs.core_properties.comments = "CONFIDENTIAL - INTERNAL USE ONLY"
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        p_color = getattr(deal, "primary_color", None) or "#002060"
        s_color = getattr(deal, "secondary_color", None) or "#800020"
        
        theme_palette = deal.theme_palette if isinstance(deal.theme_palette, list) else []
        if not theme_palette:
            theme_palette = [p_color, s_color, "#64748b", "#94a3b8", "#cbd5e1"]

        # Convert hex (e.g. "#002060") to RGBColor
        def hex_to_rgb(hex_str: str) -> RGBColor:
            if not re.fullmatch(r"#[0-9A-Fa-f]{6}", str(hex_str)):
                return RGBColor(0, 0, 0)
            hex_str = hex_str.lstrip('#')
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

        primary_rgb = hex_to_rgb(p_color)
        secondary_rgb = hex_to_rgb(s_color)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = primary_rgb

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "CREDIT PITCH BOOK"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = mask_sensitive_text(deal.customer)
        p2.font.size = Pt(28)
        p2.font.color.rgb = secondary_rgb
        p2.alignment = PP_ALIGN.CENTER

        valid_sections = [
            s for s in sorted(deal.sections, key=lambda x: x.order_index) 
            if s.state == "ready" and _section_narrative(s).strip()
        ]

        table_pattern = re.compile(r'(^[ \t]*\|[^\n]+\|[ \t]*\n[ \t]*\|[-:| ]+\|[ \t]*\n(?:[ \t]*\|[^\n]+\|[ \t]*(?:\n|$))+)', re.MULTILINE)
        sem = asyncio.Semaphore(5)

        async def process_section(section):
            async with sem:
                content = _section_narrative(section)
                charts = []
                native_tables = []
                
                # Extract markdown tables
                matches = table_pattern.findall(content)
                for table_str in matches:
                    df = _parse_table_data(table_str)
                    if df is not None:
                        try:
                            import matplotlib.pyplot as plt
                            fig, ax = plt.subplots(figsize=(10, 6))
                            colors = theme_palette if len(theme_palette) >= 5 else (theme_palette + ["#64748b", "#94a3b8", "#cbd5e1", "#334155", "#0f172a"])
                            bars = df.plot(x=df.columns[0], kind='bar', ax=ax, color=colors[:len(df.columns)-1], width=0.7)
                            
                            def format_label(val):
                                if pd.isna(val) or val == 0:
                                    return ""
                                abs_val = abs(val)
                                if abs_val >= 1_000_000:
                                    return f"{val/1_000_000:.1f}M"
                                elif abs_val >= 1_000:
                                    return f"{val/1_000:.1f}K"
                                elif abs_val >= 1:
                                    return f"{val:.1f}"
                                else:
                                    return f"{val:.2f}"
                                    
                            for container in ax.containers:
                                labels = [format_label(v.get_height()) for v in container]
                                ax.bar_label(container, labels=labels, padding=3, color='#334155', fontsize=9, fontweight='bold', rotation=90)
                                
                            ax.spines['top'].set_visible(False)
                            ax.spines['right'].set_visible(False)
                            ax.spines['left'].set_visible(False)
                            ax.set_xlabel("")
                            ax.set_yticks([])
                            plt.xticks(rotation=45, ha='right', fontsize=10)
                            ax.legend(loc='upper center', bbox_to_anchor=(0.5, 1.15), ncol=min(3, len(df.columns)-1), frameon=False)
                            plt.tight_layout()
                            
                            buf = io.BytesIO()
                            fig.savefig(buf, format='png', dpi=300, transparent=True, bbox_inches='tight')
                            plt.close(fig)
                            charts.append(buf.getvalue())
                        except Exception as e:
                            logger.error(f"Failed to generate PPT chart: {e}")
                    else:
                        lines = table_str.strip().split('\n')
                        if len(lines) >= 3:
                            headers = [col.strip() for col in lines[0].split('|')[1:-1]]
                            rows = []
                            for line in lines[2:]:
                                cols = [col.strip() for col in line.split('|')[1:-1]]
                                if len(cols) == len(headers):
                                    rows.append(cols)
                            if rows:
                                native_tables.append((headers, rows))
                
                paragraphs = [
                    re.sub(r"^[#>*\-\s]+", "", part).strip()
                    for part in re.split(r"\n\s*\n|\n(?=[*-]\s)", table_pattern.sub("", content))
                    if part.strip()
                ]
                slides_list = [{
                    "title": section.title,
                    "bullet_points": paragraphs[:6] or ["No narrative content available."],
                    "image_prompt": None,
                }]

                # Allocate generated charts to slides instead of fetching AI images
                for i, slide_info in enumerate(slides_list):
                    slide_info["img_bytes"] = None
                    if i < len(charts):
                        slide_info["img_bytes"] = charts[i]
                        slide_info["image_prompt"] = None

                return section, slides_list, native_tables

        tasks = [process_section(sec) for sec in valid_sections]
        results = await asyncio.gather(*tasks)

        for section, slides_list, native_tables in results:
            for slide_info in slides_list:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]

                title_shape.text = slide_info.get("title", section.title)
                title_shape.text_frame.paragraphs[0].font.color.rgb = primary_rgb
                
                tf = body_shape.text_frame
                tf.clear() # Clear default formatting safely
                
                for idx, point in enumerate(slide_info.get("bullet_points", [])):
                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(point)
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Gray, not white
                    p.level = 0
                    p.space_after = Pt(10)
                
                img_bytes = slide_info.get("img_bytes")
                if img_bytes:
                    try:
                        img_stream = io.BytesIO(img_bytes)
                        slide.shapes.add_picture(img_stream, Inches(7.5), Inches(2.0), width=Inches(5.0))
                        body_shape.width = Inches(6.5)
                    except Exception as e:
                        logger.warning(f"Failed to embed image: {e}")

            for headers, rows in native_tables:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"{section.title} - Data"
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_rgb
                
                sp = slide.placeholders[1]._element
                sp.getparent().remove(sp)
                
                rows_cnt = len(rows) + 1
                cols_cnt = len(headers)
                left = Inches(1.0)
                top = Inches(2.0)
                width = Inches(11.0)
                height = Inches(0.5 * rows_cnt)
                
                table_shape = slide.shapes.add_table(rows_cnt, cols_cnt, left, top, width, height)
                table = table_shape.table
                
                for col_idx, header in enumerate(headers):
                    if col_idx >= cols_cnt: break
                    cell = table.cell(0, col_idx)
                    cell.text = header
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = primary_rgb
                    
                for row_idx, row_data in enumerate(rows):
                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx < cols_cnt:
                            cell = table.cell(row_idx + 1, col_idx)
                            cell.text = str(cell_data)
                            cell.text_frame.paragraphs[0].font.size = Pt(12)
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    @staticmethod
    def generate_combined_report(deal: Deal) -> bytes:
        return ExportService.generate_pdf(deal)
